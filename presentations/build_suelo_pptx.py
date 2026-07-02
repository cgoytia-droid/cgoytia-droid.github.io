#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Genera "Luján de Cuyo · Innovación en gestión del suelo" como PowerPoint editable (.pptx).
Rediseño académico y motivador, de bajo texto: máximo 3 bullets por slide, tipografía
Calibri de buen tamaño. Misma narrativa que el material fuente (CIPUV-UTDT / Lincoln Institute).

Uso:  python3 build_suelo_pptx.py
Salida: lujan-de-cuyo-gestion-suelo.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- paleta (sitio de la autora + acentos suelo/agua) ----------
SLATE   = RGBColor(0x4E,0x58,0x69)
SLATE2  = RGBColor(0x6F,0x7D,0x94)
SLATE_DK= RGBColor(0x3A,0x42,0x50)
GOLD    = RGBColor(0xB8,0x86,0x0B)
GOLD_LT = RGBColor(0xD4,0xC0,0x88)
INK     = RGBColor(0x2B,0x30,0x38)
PAPER   = RGBColor(0xF7,0xF8,0xFA)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
MUTED   = RGBColor(0x7A,0x82,0x8E)
WATER   = RGBColor(0x2F,0x8F,0x9D)   # agua / azul-verde
WATER_BG= RGBColor(0xE7,0xF1,0xF3)
GREEN   = RGBColor(0x5A,0x8F,0x6B)   # verde / recarga
GREEN_BG= RGBColor(0xE7,0xF0,0xEA)
CLAY    = RGBColor(0xB1,0x45,0x6A)   # costo / sello del suelo
CLAY_BG = RGBColor(0xF1,0xE6,0xEA)
CARD_BD = RGBColor(0xE3,0xE6,0xEA)

FONT = "Calibri"
FONT_LT = "Calibri Light"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

import os
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def _set_font(run, size, color, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = name; run.font.color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, round_=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    if round_ is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = round_
        except Exception: pass
    return sp

def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
        {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val':'28323E'})
    alpha = clr.makeelement(qn('a:alpha'), {'val':'20000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=8, line_spacing=1.06, wrap=True):
    """runs: lista de párrafos; cada párrafo es lista de (txt,size,color,bold,italic[,name])."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            seg = list(seg)
            txt, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            italic = seg[4] if len(seg) > 4 else False
            name = seg[5] if len(seg) > 5 else FONT
            r = p.add_run(); r.text = txt; _set_font(r, size, color, bold, italic, name)
    return tb

def base(s, dark=False, bg=None):
    b = s.background; b.fill.solid()
    b.fill.fore_color.rgb = bg if bg is not None else (SLATE_DK if dark else PAPER)

def brandbar(s):
    rect(s, 0, 0, EMU_W, Pt(7), fill=GOLD)

def kicker(s, label, x=Inches(0.75), y=Inches(0.62), color=GOLD):
    rect(s, x, y+Emu(24000), Pt(5), Inches(0.30), fill=color)
    text(s, x+Inches(0.16), y, Inches(11), Inches(0.36),
         [[(label.upper(), 15, color, True, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)

def title(s, parts, x=Inches(0.75), y=Inches(1.08), w=Inches(11.8), size=36):
    runs=[[(t, size, c, True, it) for (t,c,it) in parts]]
    text(s, x, y, w, Inches(1.3), runs, line_spacing=1.0, space_after=0)

TOTAL = 29  # total de slides (numeración automática)
def pagenum(s):
    n = len(prs.slides._sldIdLst)  # índice del slide actual (recién creado)
    text(s, Inches(12.2), Inches(6.92), Inches(0.9), Inches(0.4),
         [[(f"{n:02d} / {TOTAL}", 12, MUTED, False, False)]], align=PP_ALIGN.RIGHT,
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)

def footer(s):
    text(s, Inches(0.75), Inches(6.92), Inches(9), Inches(0.4),
         [[("Luján de Cuyo · Gestión del suelo   ·   ", 11, MUTED, False, False),
           ("CIPUV–UTDT / Lincoln Institute", 11, SLATE2, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)

def bullets(s, items, x=Inches(0.75), y=Inches(2.35), w=Inches(8.4), size=20,
            color=INK, gap=16, dot=GOLD, lh=1.12):
    """items: lista de strings o de listas de runs (para negritas). Máx 3 recomendado."""
    tb = s.shapes.add_textbox(x, y, w, Inches(4)); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = lh
        r = p.add_run(); r.text = "▪  "; _set_font(r, size, dot, True, False)
        if isinstance(it, str):
            it = [(it, size, color, False, False)]
        for seg in it:
            seg = list(seg) + [False, False]
            t, sz, c, b, ital = seg[:5]
            r = p.add_run(); r.text = t; _set_font(r, sz, c, b, ital)
    return tb

def stat_card(s, x, y, w, h, big, small, accent=GOLD, big_size=40, small_size=13.5,
              bg=WHITE):
    rect(s, x, y, w, h, fill=bg, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
    rect(s, x, y, w, Pt(5), fill=accent)
    text(s, x+Inches(0.22), y+Inches(0.24), w-Inches(0.44), h-Inches(0.4),
         [[(big, big_size, accent, True, False)],
          [(small, small_size, SLATE, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=6, line_spacing=1.05)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def chip(s, x, y, w, label, color=GOLD, txtcolor=WHITE, h=Inches(0.42), size=13):
    rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
    text(s, x, y, w, h, [[(label, size, txtcolor, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

# =====================================================================
# SLIDE 1 — PORTADA
# =====================================================================
s = slide(); base(s, bg=SLATE_DK)
# bloque lateral dorado
rect(s, 0, 0, Inches(0.32), EMU_H, fill=GOLD)
# franja de acento agua->verde
rect(s, Inches(0.32), Inches(6.9), EMU_W-Inches(0.32), Inches(0.6), fill=WATER)
rect(s, Inches(6.8), Inches(6.9), EMU_W-Inches(6.8), Inches(0.6), fill=GREEN)
text(s, Inches(0.95), Inches(0.85), Inches(11), Inches(0.4),
     [[("CIPUV–UTDT · LINCOLN INSTITUTE", 15, GOLD_LT, True, False)]], space_after=0)
text(s, Inches(0.95), Inches(1.35), Inches(11.5), Inches(0.5),
     [[("LUJÁN DE CUYO · MENDOZA, ARGENTINA", 14, RGBColor(0xC9,0xD0,0xDA), True, False)]], space_after=0)
text(s, Inches(0.95), Inches(2.15), Inches(11.4), Inches(2.2),
     [[("Innovación en gestión del suelo", 50, WHITE, True, False, FONT)],
      [("para financiar la adaptación climática", 50, GOLD_LT, True, False, FONT)]],
     line_spacing=1.02, space_after=2)
text(s, Inches(0.95), Inches(4.55), Inches(11.2), Inches(0.6),
     [[("Seis instrumentos de captura de valor · un Fideicomiso · el Índice de Ciudad Deseada",
        20, RGBColor(0xD8,0xDD,0xE4), False, True)]], space_after=0)
# meta chips
chip(s, Inches(0.95), Inches(5.5), Inches(3.2), "Caso: Mendoza, Argentina", color=WATER)
chip(s, Inches(4.35), Inches(5.5), Inches(3.5), "Enfoque: Land Value Capture + clima", color=GREEN)
chip(s, Inches(8.05), Inches(5.5), Inches(3.0), "Clase · 30 minutos", color=GOLD)
notes(s, "Bienvenida. Un caso real de una ciudad intermedia que usa la gestión del suelo "
          "para financiar su propia adaptación climática. Hoy: seis instrumentos, un fideicomiso "
          "y una regla de gasto con equidad.")

# =====================================================================
# SLIDE 2 — HOJA DE RUTA
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Hoja de ruta")
title(s, [("Lo que veremos en 30 minutos", SLATE, False)])
road = [
    ("01", "El contexto y la paradoja agua–suelo", WATER),
    ("02", "La idea: captura de valor del suelo", GOLD),
    ("03", "Los seis instrumentos", CLAY),
    ("04", "La arquitectura: Banco + Fondo + Fideicomiso", SLATE),
    ("05", "La regla de gasto: Índice de Ciudad Deseada", GREEN),
    ("06", "Números, equidad y co-beneficios", GOLD),
]
cw, ch, gx, gy = Inches(5.9), Inches(1.15), Inches(0.28), Inches(0.28)
x0, y0 = Inches(0.75), Inches(2.35)
for i,(n,t,c) in enumerate(road):
    col, row = i % 2, i // 2
    x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
    rect(s, x, y, cw, ch, fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
    rect(s, x, y, Inches(0.12), ch, fill=c)
    text(s, x+Inches(0.32), y, Inches(1.1), ch, [[(n, 34, c, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x+Inches(1.45), y, cw-Inches(1.65), ch, [[(t, 19, INK, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.05)
footer(s); pagenum(s)
notes(s, "El arco de la clase: del problema físico (agua-suelo) a una máquina financiera "
          "que captura valor, lo blinda y lo redistribuye con equidad.")

# =====================================================================
# SLIDE 3 — LA IDEA EN UNA FRASE
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "La idea en una frase")
title(s, [("Cuatro verbos, un sistema", SLATE, False)])
ideas = [
    ("Capturar", "la plusvalía que genera el propio desarrollo urbano", WATER, WATER_BG),
    ("Blindar", "los recursos en un fideicomiso", SLATE, RGBColor(0xEC,0xEE,0xF1)),
    ("Apalancar", "financiamiento climático internacional", GOLD, RGBColor(0xF6,0xF0,0xE0)),
    ("Redistribuir", "con equidad, hacia donde más se necesita", GREEN, GREEN_BG),
]
cw2, ch2, gx2 = Inches(2.85), Inches(3.1), Inches(0.28)
x0 = Inches(0.75); y0 = Inches(2.55)
for i,(h,d,c,bg) in enumerate(ideas):
    x = x0 + i*(cw2+gx2)
    rect(s, x, y0, cw2, ch2, fill=bg, line=None, shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
    rect(s, x, y0, cw2, Pt(6), fill=c)
    text(s, x+Inches(0.28), y0+Inches(0.55), cw2-Inches(0.5), Inches(0.7),
         [[(f"{i+1}", 22, c, True, False)]], space_after=0)
    text(s, x+Inches(0.28), y0+Inches(1.0), cw2-Inches(0.5), Inches(0.7),
         [[(h, 27, c, True, False)]], space_after=0)
    text(s, x+Inches(0.28), y0+Inches(1.75), cw2-Inches(0.5), Inches(1.2),
         [[(d, 17, INK, False, False)]], space_after=0, line_spacing=1.12)
footer(s); pagenum(s)
notes(s, "Toda la clase cabe en cuatro verbos: capturar, blindar, apalancar, redistribuir. "
          "Si se recuerda esto, se recuerda el caso.")

# =====================================================================
# SLIDE 4 — EL CONTEXTO
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "El contexto · 01", color=WATER)
title(s, [("Una ciudad que crece con ", SLATE, False), ("presupuesto de agua", WATER, False)])
bullets(s, [
    [("~175.000 habitantes", 21, INK, True), (" en un oasis semiárido al pie de los Andes.", 21, INK, False)],
    [("~200 mm de lluvia al año:", 21, INK, True), (" todo el sistema depende del deshielo andino.", 21, INK, False)],
    [("Dos tendencias en colisión:", 21, INK, True), (" menos nieve, más urbanización del periurbano regado.", 21, INK, False)],
], w=Inches(8.1), dot=WATER, gap=20)
stat_card(s, Inches(9.35), Inches(2.45), Inches(3.1), Inches(2.9),
          "+46%", "crecimiento poblacional 2010–2022", accent=WATER, big_size=52, small_size=15)
footer(s); pagenum(s)
notes(s, "El dato duro: crece 46% en doce años sobre 200 mm de lluvia. Cada hectárea que se "
          "urbaniza suma demanda de agua y resta recarga del acuífero.")

# =====================================================================
# SLIDE (nuevo) — POBLACIÓN: MENDOZA Y LUJÁN DE CUYO (mapa real CIPUV)
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "El contexto · 01", color=WATER)
title(s, [("¿Dónde crece Luján de Cuyo?", SLATE, False)])
# mapa (choropleth por distrito, 2022) a la derecha
_mapw = Inches(6.9); _maph = Emu(int(_mapw * 540 / 764))
s.shapes.add_picture(os.path.join(ASSETS, "mapa-poblacion-lujan.png"),
                     Inches(6.25), Inches(1.95), width=_mapw, height=_maph)
text(s, Inches(6.25), Inches(1.95)+_maph+Emu(20000), _mapw, Inches(0.3),
     [[("Población por distrito, 2022 · Gran Mendoza y Luján de Cuyo", 11, MUTED, False, True)]],
     space_after=0)
# bullets a la izquierda
bullets(s, [
    [("Carrodilla (30.469) y Perdriel (24.922)", 19, INK, True), (" concentran el mayor tamaño y crecimiento.", 19, INK, False)],
    [("El avance urbano ocurre sobre el ", 19, INK, False), ("periurbano regado", 19, WATER, True), (": Chacras, Vistalba, Vertientes.", 19, INK, False)],
    [("Cada distrito que crece ", 19, INK, False), ("urbaniza suelo con derecho de agua", 19, GREEN, True), (" → más demanda, menos recarga.", 19, INK, False)],
], x=Inches(0.75), y=Inches(2.15), w=Inches(5.25), dot=WATER, gap=18, size=19)
stat_card(s, Inches(0.75), Inches(5.15), Inches(5.25), Inches(1.35),
          "+42,6%  ·  +50.546 hab.", "crecimiento intercensal 2010–2022 en Luján de Cuyo",
          accent=WATER, big_size=24, small_size=14)
footer(s); pagenum(s)
notes(s, "El mapa de la propia CIPUV: el crecimiento no es uniforme. Se concentra en el este, "
          "sobre el oasis regado del Gran Mendoza —justo el suelo cuya conversión dispara la paradoja "
          "agua–suelo. Fuente: elaboración propia en base a Municipalidad de Luján de Cuyo e INDEC.")

# =====================================================================
# SLIDE 5 — LA PARADOJA AGUA–SUELO
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "El problema · 01", color=CLAY)
title(s, [("La paradoja agua–suelo", SLATE, False)])
# dos cajas
bx, by, bw, bh = Inches(0.75), Inches(2.4), Inches(5.7), Inches(2.5)
rect(s, bx, by, bw, bh, fill=WATER_BG, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
text(s, bx+Inches(0.3), by+Inches(0.25), bw-Inches(0.6), Inches(0.5),
     [[("GANANCIA PRIVADA", 16, WATER, True, False)]], space_after=0)
text(s, bx+Inches(0.3), by+Inches(0.85), bw-Inches(0.6), Inches(1.5),
     [[("Rezonificar un viñedo a loteo multiplica su valor varias veces.", 20, INK, False, False)]],
     space_after=0, line_spacing=1.15)
bx2 = Inches(6.85)
rect(s, bx2, by, bw, bh, fill=CLAY_BG, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
text(s, bx2+Inches(0.3), by+Inches(0.25), bw-Inches(0.6), Inches(0.5),
     [[("COSTO COLECTIVO", 16, CLAY, True, False)]], space_after=0)
text(s, bx2+Inches(0.3), by+Inches(0.85), bw-Inches(0.6), Inches(1.5),
     [[("El mismo acto sella el suelo, elimina recarga, sube la demanda de agua y borra el sumidero de carbono.", 20, INK, False, False)]],
     space_after=0, line_spacing=1.15)
# franja conclusión
rect(s, Inches(0.75), Inches(5.25), Inches(11.8), Inches(1.15), fill=SLATE_DK,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
text(s, Inches(1.1), Inches(5.25), Inches(11.1), Inches(1.15),
     [[("→  El incremento de valor es la base natural para financiar la adaptación: ", 20, WHITE, False, False),
       ("quien causa el costo, financia la respuesta.", 20, GOLD_LT, True, False)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.1)
footer(s); pagenum(s)
notes(s, "El corazón del caso. Una sola decisión —que otorga el propio Concejo— produce ganancia "
          "privada y costo colectivo. Ese diferencial de valor es el que puede pagar la adaptación.")

# =====================================================================
# SLIDE 6 — EL MARCO (LVC)
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "El marco · 02", color=GOLD)
title(s, [("¿Qué es la captura de valor del suelo?", SLATE, False)])
bullets(s, [
    [("El precio del suelo capitaliza la inversión pública, la regulación y el crecimiento.", 21, INK, False)],
    [("Buena parte de ese incremento es ", 21, INK, False), ("«no ganado»", 21, GOLD, True),
     (" por el propietario (George; Glaeser; Duranton & Puga).", 21, INK, False)],
    [("Gravarlo es ", 21, INK, False), ("eficiente y justo", 21, GREEN, True),
     (": recae sobre una ganancia que no produjo quien la recibe.", 21, INK, False)],
], w=Inches(11.6), gap=20)
text(s, Inches(0.75), Inches(5.9), Inches(11.7), Inches(0.9),
     [[("Tradición latinoamericana: valorización (Colombia), CEPACs (São Paulo), Estatuto de la Ciudad (Brasil). ",
        15, MUTED, False, True),
       ("La brecha que este caso cierra: unir la LVC con el financiamiento climático (Dunning & Lord, 2020).",
        15, SLATE2, False, True)]], space_after=0, line_spacing=1.12)
footer(s); pagenum(s)
notes(s, "El fundamento teórico: la plusvalía del suelo es en buena parte no ganada. Gravarla es "
          "eficiente (no distorsiona) y justo. La novedad es acoplarla al financiamiento climático.")

# =====================================================================
# SLIDE 7 — POR QUÉ IMPORTA
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Por qué importa · 02", color=GOLD)
title(s, [("El dinero sigue la ", SLATE, False), ("capacidad", CLAY, False),
          (", no la ", SLATE, False), ("necesidad", GREEN, False)])
stat_card(s, Inches(0.75), Inches(2.45), Inches(3.6), Inches(3.4),
          "USD 147\nmil M/año", "de adaptación que necesitan las ciudades en desarrollo a 2030 (CCFLA, 2024)",
          accent=CLAY, big_size=34, small_size=15)
bullets(s, [
    [("El financiamiento fluye a las metrópolis con más capacidad, no a las más expuestas (Clifton et al., 2025).", 20, INK, False)],
    [("Las ciudades intermedias quedan atrapadas: ", 20, INK, False), ("sin capacidad no acceden; sin acceso no construyen capacidad.", 20, CLAY, True)],
    [("La salida: ", 20, INK, False), ("generar recurso propio, predecible y blindado", 20, GREEN, True), (" — la «llave» que pide el financista.", 20, INK, False)],
], x=Inches(4.75), y=Inches(2.5), w=Inches(7.7), gap=22)
footer(s); pagenum(s)
notes(s, "El problema estructural: el capital climático premia a quien ya puede, no a quien más "
          "lo necesita. La única salida para una ciudad intermedia es traer su propia llave.")

# =====================================================================
# SLIDE 8 — LA INNOVACIÓN (SISTEMA)
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "La innovación · 02")
title(s, [("Un sistema, no instrumentos sueltos", SLATE, False)])
steps = [
    ("Crear valor", "obra · rezonificación", WATER),
    ("Capturar", "6 instrumentos LVC", GOLD),
    ("Blindar", "Banco + Fondo + Fideicomiso", SLATE),
    ("Redistribuir", "Índice de Ciudad Deseada", GREEN),
    ("Apalancar", "bono verde · multilaterales", CLAY),
]
n = len(steps); sw = Inches(2.15); gap = Inches(0.20)
x0 = Inches(0.75); y0 = Inches(2.9); sh = Inches(2.0)
for i,(h,d,c) in enumerate(steps):
    x = x0 + i*(sw+gap)
    rect(s, x, y0, sw, sh, fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.07)
    rect(s, x, y0, sw, Pt(6), fill=c)
    text(s, x+Inches(0.18), y0+Inches(0.28), sw-Inches(0.36), Inches(0.5),
         [[(f"0{i+1}", 20, c, True, False)]], space_after=0)
    text(s, x+Inches(0.18), y0+Inches(0.78), sw-Inches(0.36), Inches(0.55),
         [[(h, 19, INK, True, False)]], space_after=0)
    text(s, x+Inches(0.18), y0+Inches(1.28), sw-Inches(0.36), Inches(0.65),
         [[(d, 13.5, SLATE2, False, False)]], space_after=0, line_spacing=1.05)
    if i < n-1:
        text(s, x+sw-Inches(0.02), y0, gap+Inches(0.1), sh,
             [[("→", 20, GOLD, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, Inches(0.75), Inches(5.35), Inches(11.8), Inches(0.8),
     [[("La coordinación es lo que convierte 6 cargos modestos en un ", 20, INK, False, False),
       ("sistema bancable", 20, GOLD, True, False), (".", 20, INK, False, False)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=0)
footer(s); pagenum(s)
notes(s, "La tesis central: el valor no está en cada instrumento, sino en encadenarlos. "
          "Crear → capturar → blindar → redistribuir → apalancar.")

# =====================================================================
# SLIDE 9 — LOS SEIS INSTRUMENTOS (divisor)
# =====================================================================
s = slide(); base(s, bg=SLATE_DK)
rect(s, 0, 0, Inches(0.32), EMU_H, fill=GOLD)
text(s, Inches(0.95), Inches(1.1), Inches(11), Inches(0.5),
     [[("SECCIÓN 03", 16, GOLD_LT, True, False)]], space_after=0)
text(s, Inches(0.95), Inches(1.7), Inches(11.5), Inches(1.2),
     [[("Los seis instrumentos", 44, WHITE, True, False)]], space_after=0)
text(s, Inches(0.95), Inches(2.85), Inches(11), Inches(0.7),
     [[("Cada uno captura el incremento de valor en un punto distinto del ciclo de desarrollo.",
        20, RGBColor(0xCF,0xD5,0xDE), False, True)]], space_after=0, line_spacing=1.15)
insts = [
    ("1", "Contribución por Mejoras", WATER),
    ("2", "Participación en Valoración", WATER),
    ("3", "Sobretasa a baldíos", GREEN),
    ("4", "Suelo Creado / DCA", GOLD),
    ("5", "Incentivos fiscales", GOLD),
    ("6", "Transferencia de Derechos", CLAY),
]
cw, ch, gx, gy = Inches(3.75), Inches(0.95), Inches(0.28), Inches(0.26)
x0, y0 = Inches(0.95), Inches(3.9)
for i,(n,t,c) in enumerate(insts):
    col, row = i % 3, i // 3
    x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
    rect(s, x, y, cw, ch, fill=RGBColor(0x45,0x4E,0x5C), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
    rect(s, x, y, Inches(0.7), ch, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
    text(s, x, y, Inches(0.7), ch, [[(n, 24, WHITE, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x+Inches(0.85), y, cw-Inches(1.0), ch, [[(t, 17, WHITE, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.05)
notes(s, "Vista general de los seis. No hace falta memorizarlos: lo importante es que cada uno "
          "engancha en un momento distinto del ciclo del desarrollo.")

# =====================================================================
# Helper para slides de instrumento
# =====================================================================
def instrumento(n_page, kick, ttl, num, estado, estado_color, items, extra_stats=None, accent=WATER):
    s = slide(); base(s); brandbar(s); kicker(s, kick, color=accent)
    title(s, [(ttl, SLATE, False)])
    # número grande a la izquierda
    rect(s, Inches(0.75), Inches(2.35), Inches(1.5), Inches(1.5), fill=accent,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.14, shadow=True)
    text(s, Inches(0.75), Inches(2.35), Inches(1.5), Inches(1.5),
         [[(num, 60, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    chip(s, Inches(0.75), Inches(4.05), Inches(1.5), estado, color=estado_color, size=12)
    bullets(s, items, x=Inches(2.6), y=Inches(2.35),
            w=Inches(6.6) if extra_stats else Inches(9.9), dot=accent, gap=18, size=20)
    if extra_stats:
        yy = Inches(2.45)
        for (big, small, c) in extra_stats:
            stat_card(s, Inches(9.5), yy, Inches(2.95), Inches(1.35), big, small,
                      accent=c, big_size=26, small_size=13)
            yy += Inches(1.55)
    footer(s); pagenum(s)  # n_page ignorado: numeración automática
    return s

# SLIDE 10 — Instrumento 1
s = instrumento(10, "Instrumento 1 · 03", "Contribución por Mejoras (betterment)", "1",
    "DISEÑADO", SLATE2, [
    [("Grava el aumento de valor que una ", 20, INK, False), ("obra pública", 20, WATER, True),
     (" genera en los inmuebles vecinos.", 20, INK, False)],
    [("Tasa: ", 20, INK, False), ("30–50% del incremento", 20, GOLD, True),
     (" en una «zona de beneficio» delimitada por SIG.", 20, INK, False)],
    [("Doble función: ", 20, INK, False), ("financia la obra hídrica y demuestra, en efectivo, que rindió.", 20, GREEN, True)],
], accent=WATER)
notes(s, "Instrumento clásico. Si el municipio hace la obra hídrica, captura parte de la valorización "
          "que esa obra crea. Se paga en cuotas para no golpear a los hogares de menores ingresos.")

# SLIDE 11 — Instrumento 2
s = instrumento(11, "Instrumento 2 · 03", "Participación en la Valoración de Activos", "2",
    "DISEÑADO", SLATE2, [
    [("Grava la plusvalía por ", 20, INK, False), ("rezonificación o mayor intensidad de uso.", 20, WATER, True)],
    [("Tasa: ", 20, INK, False), ("30–50% del diferencial", 20, GOLD, True), (" de valor antes / después.", 20, INK, False)],
    [("Es la lógica del sistema en miniatura: ", 20, INK, False),
     ("el acto que genera el costo es el que paga la respuesta.", 20, GREEN, True)],
], accent=WATER)
notes(s, "Este es el instrumento emblemático de la paradoja: la conversión de suelo permeable a "
          "impermeable genera el costo, y este mismo cargo captura el pago que lo financia.")

# SLIDE 12 — Instrumento 3
s = instrumento(12, "Instrumento 3 · 03", "Sobretasa a Inmueble Ocioso (baldíos)", "3",
    "● OPERATIVO", GREEN, [
    [("Grava parcelas ", 20, INK, False), ("servidas pero ociosas", 20, GREEN, True), (" dentro del perímetro urbano.", 20, INK, False)],
    [("Recargo progresivo: ", 20, INK, False), ("+100%", 20, GOLD, True), (" estándar y ", 20, INK, False), ("+200%", 20, CLAY, True), (" en vacíos prioritarios.", 20, INK, False)],
    [("Promueve el ", 20, INK, False), ("desarrollo compacto", 20, GREEN, True), (" usando el padrón que ya existe.", 20, INK, False)],
], extra_stats=[
    ("5.315", "lotes alcanzados (2026)", GREEN),
    ("~ARS 345 M", "/año  ≈ USD 277.500", GOLD),
], accent=GREEN)
notes(s, "Ya está operativo. 5.315 lotes, recaudación real estimada. Además del ingreso, empuja el "
          "relleno urbano en vez del sprawl: adaptación y forma urbana a la vez.")

# SLIDE 13 — Instrumento 4
s = instrumento(13, "Instrumento 4 · 03", "Suelo Creado / Derechos de Construcción Adicional", "4",
    "DISEÑADO", SLATE2, [
    [("Vende el ", 20, INK, False), ("derecho a construir por encima de una «línea de base».", 20, WATER, True)],
    [("Modelo ", 20, INK, False), ("CEPAC (São Paulo)", 20, GOLD, True), (" adaptado a ciudad intermedia.", 20, INK, False)],
    [("COS 2026: de 10 m a ", 20, INK, False), ("36 m (~8,7 pisos extra)", 20, CLAY, True), (" sólo vía compra; bonos por eficiencia hídrica.", 20, INK, False)],
], accent=GOLD)
notes(s, "La edificabilidad extra deja de ser un regalo del zoning y pasa a comprarse. El precio se "
          "referencia al mercado y premia el diseño eficiente en agua.")

# SLIDE 14 — Instrumentos 5 y 6
s = slide(); base(s); brandbar(s); kicker(s, "Instrumentos 5 y 6 · 03", color=GOLD)
title(s, [("Incentivos fiscales · Transferencia de Derechos", SLATE, False)])
for i,(num,h,body,c,bg) in enumerate([
    ("5","Incentivos fiscales","Un «gradiente verde»: descuentos por eficiencia hídrica, permeabilidad y arbolado. Sin adaptación, cargo pleno; con adaptación, alivio fiscal.",GREEN,GREEN_BG),
    ("6","Transferencia de Derechos (TDR)","Proteger recarga, riberas y agro sin comprar el suelo: las zonas emisoras venden sus derechos a zonas receptoras vía mercado.",WATER,WATER_BG),
]):
    x = Inches(0.75) + i*Inches(6.1)
    rect(s, x, Inches(2.4), Inches(5.8), Inches(3.6), fill=bg, shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
    rect(s, x, Inches(2.4), Inches(5.8), Pt(6), fill=c)
    rect(s, x+Inches(0.35), Inches(2.75), Inches(0.9), Inches(0.9), fill=c,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.16)
    text(s, x+Inches(0.35), Inches(2.75), Inches(0.9), Inches(0.9),
         [[(num, 36, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x+Inches(1.45), Inches(2.85), Inches(4.1), Inches(0.9),
         [[(h, 22, c, True, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.02)
    text(s, x+Inches(0.4), Inches(4.0), Inches(5.05), Inches(1.85),
         [[(body, 18, INK, False, False)]], space_after=0, line_spacing=1.2)
footer(s); pagenum(s)
notes(s, "Dos instrumentos que orientan conducta, no sólo recaudan: el gradiente verde premia el "
          "diseño adaptado, y la TDR protege suelo estratégico sin gasto de compra.")

# =====================================================================
# SLIDE 15 — INSTRUMENTO ESTRELLA
# =====================================================================
s = slide(); base(s, bg=SLATE_DK)
rect(s, 0, 0, Inches(0.32), EMU_H, fill=GOLD)
text(s, Inches(0.95), Inches(0.75), Inches(11), Inches(0.4),
     [[("(+) EL INSTRUMENTO ESTRELLA", 15, GOLD_LT, True, False)]], space_after=0)
text(s, Inches(0.95), Inches(1.25), Inches(11.5), Inches(0.9),
     [[("Aporte de Agua y Cloaca — la evidencia más fuerte", 32, WHITE, True, False)]],
     space_after=0, line_spacing=1.02)
chip(s, Inches(0.95), Inches(2.25), Inches(1.6), "● OPERATIVO", color=GREEN, size=13)
# progresión 2021 -> 2025
rect(s, Inches(0.95), Inches(3.15), Inches(3.4), Inches(2.3), fill=RGBColor(0x45,0x4E,0x5C),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
text(s, Inches(0.95), Inches(3.35), Inches(3.4), Inches(2.0),
     [[("2021", 18, RGBColor(0xB9,0xC0,0xCA), True, False)],
      [("USD 60.700", 30, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=8)
text(s, Inches(4.5), Inches(3.15), Inches(1.2), Inches(2.3),
     [[("→", 44, GOLD, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
rect(s, Inches(5.85), Inches(3.15), Inches(3.4), Inches(2.3), fill=GREEN,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06, shadow=True)
text(s, Inches(5.85), Inches(3.35), Inches(3.4), Inches(2.0),
     [[("2025", 18, RGBColor(0xDD,0xEA,0xE1), True, False)],
      [("USD 690.500", 30, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=8)
# 11x
rect(s, Inches(9.7), Inches(3.15), Inches(2.75), Inches(2.3), fill=GOLD,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08, shadow=True)
text(s, Inches(9.7), Inches(3.15), Inches(2.75), Inches(2.3),
     [[("11×", 60, WHITE, True, False)],[("en cuatro años", 15, WHITE, False, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=4)
text(s, Inches(0.95), Inches(5.8), Inches(11.5), Inches(0.9),
     [[("Prueba de concepto, en efectivo: un cargo ligado al suelo escala en pocos años.",
        20, RGBColor(0xCF,0xD5,0xDE), False, True)]], space_after=0, line_spacing=1.1)
notes(s, "La evidencia más contundente. No es teoría: un cargo por conexión de servicios se multiplicó "
          "por 11 en cuatro años, en dólares constantes. Es la prueba de que el modelo escala.")

# =====================================================================
# SLIDE 16 — SÍNTESIS: TABLA
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Síntesis")
title(s, [("Los instrumentos de un vistazo", SLATE, False)])
rows = [
    ("Contribución por Mejoras", "Obra pública", "30–50% del incremento", "DISEÑADO", SLATE2),
    ("Participación en Valoración", "Rezonificación", "30–50% del diferencial", "DISEÑADO", SLATE2),
    ("Sobretasa a baldíos", "Lote ocioso servido", "+100% / +200%", "● OPERATIVO", GREEN),
    ("Suelo Creado / DCA", "Edificar sobre la base", "Mercado (COS 2026)", "DISEÑADO", SLATE2),
    ("Incentivos fiscales", "Diseño climático", "Alivio variable", "DISEÑADO", SLATE2),
    ("Transferencia de Derechos", "Zona receptora", "Mercado", "DISEÑADO", SLATE2),
    ("Aporte Agua y Cloaca", "Nuevo desarrollo", "Tarifaria", "● OPER · 11×", GREEN),
]
cols = [Inches(4.2), Inches(3.1), Inches(2.9), Inches(1.6)]
heads = ["INSTRUMENTO", "DISPARADOR", "TASA / BASE", "ESTADO"]
x0, y0 = Inches(0.75), Inches(2.25); rh = Inches(0.58)
# header
xx = x0
rect(s, x0, y0, sum(cols, Emu(0)), rh, fill=SLATE_DK)
for cwid, hd in zip(cols, heads):
    text(s, xx+Inches(0.15), y0, cwid-Inches(0.2), rh, [[(hd, 13.5, WHITE, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    xx += cwid
# rows
for i,(a,b,c,st,stc) in enumerate(rows):
    y = y0 + rh + i*Inches(0.55)
    if i % 2 == 0:
        rect(s, x0, y, sum(cols, Emu(0)), Inches(0.55), fill=WHITE)
    else:
        rect(s, x0, y, sum(cols, Emu(0)), Inches(0.55), fill=RGBColor(0xEF,0xF1,0xF4))
    xx = x0
    vals = [(a, INK, True), (b, SLATE, False), (c, SLATE, False), (st, stc, True)]
    for cwid,(v,vc,vb) in zip(cols, vals):
        text(s, xx+Inches(0.15), y, cwid-Inches(0.2), Inches(0.55),
             [[(v, 14, vc, vb, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        xx += cwid
footer(s); pagenum(s)
notes(s, "Cuadro de referencia. Dos ya operativos (verde), el resto diseñados y listos. "
          "Nótese la diversidad de disparadores: cada uno cubre un momento distinto.")

# =====================================================================
# SLIDE 17 — ¿POR QUÉ UN FIDEICOMISO?
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "La arquitectura · 04", color=SLATE)
title(s, [("¿Por qué un ", SLATE, False), ("Fideicomiso", GOLD, False), ("?", SLATE, False)])
bullets(s, [
    [("Recurso ", 20, INK, False), ("blindado", 20, GOLD, True), (", legalmente separado del presupuesto general.", 20, INK, False)],
    [("Es el ", 20, INK, False), ("«colateral»", 20, WATER, True), (" que un banco o un bono puede ver y auditar.", 20, INK, False)],
    [("Gobernanza que ", 20, INK, False), ("blinda los recursos entre gestiones", 20, GREEN, True), (": directorio + externos, auditoría publicada.", 20, INK, False)],
], x=Inches(0.75), y=Inches(2.4), w=Inches(7.4), gap=20)
# reparto 60/25/15
alloc = [("60%","Inversión en infraestructura",WATER),("25%","Operación y mantenimiento",SLATE2),("15%","Cuenta de Apalancamiento",GOLD)]
yy = Inches(2.45)
for big, small, c in alloc:
    rect(s, Inches(8.5), yy, Inches(3.95), Inches(1.1), fill=WHITE, line=CARD_BD, line_w=Pt(1),
         shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
    rect(s, Inches(8.5), yy, Inches(0.14), Inches(1.1), fill=c)
    text(s, Inches(8.8), yy, Inches(1.5), Inches(1.1), [[(big, 30, c, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, Inches(10.25), yy, Inches(2.1), Inches(1.1), [[(small, 14, SLATE, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.05)
    yy += Inches(1.25)
footer(s); pagenum(s)
notes(s, "El fideicomiso es la pieza que convierte recaudación en credibilidad. Separado del "
          "presupuesto, auditado, con reglas de reparto fijas: eso es lo que un financista puede confiar.")

# =====================================================================
# SLIDE 18 — ESCALERA DE APALANCAMIENTO
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "De recurso propio a capital global · 04", color=GOLD)
title(s, [("La escalera de apalancamiento", SLATE, False)])
ladder = [
    ("01","Matching","La Cuenta de Apalancamiento (~USD 0,45 M/año) ancla la contrapartida de 10–30%.",WATER),
    ("02","Emitir","Bono verde de 3–4× la recaudación anual (~USD 9–12 M).",GOLD),
    ("03","Escalar","El historial abre facilidades mayores: CAF, BID, GCF.",GREEN),
]
for i,(n,h,body,c) in enumerate(ladder):
    y = Inches(2.5) + i*Inches(1.25)
    x = Inches(0.75) + i*Inches(1.1)   # efecto escalera
    rect(s, x, y, Inches(9.8), Inches(1.05), fill=WHITE, line=CARD_BD, line_w=Pt(1),
         shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
    rect(s, x, y, Inches(1.0), Inches(1.05), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
    text(s, x, y, Inches(1.0), Inches(1.05), [[(n, 26, WHITE, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x+Inches(1.25), y+Inches(0.12), Inches(8.3), Inches(0.45), [[(h, 20, c, True, False)]], space_after=0)
    text(s, x+Inches(1.25), y+Inches(0.55), Inches(8.3), Inches(0.45), [[(body, 15.5, INK, False, False)]], space_after=0, line_spacing=1.05)
text(s, Inches(0.75), Inches(6.5), Inches(11.6), Inches(0.4),
     [[("La credibilidad, como el capital, se acumula.", 18, GOLD, True, True)]], space_after=0)
footer(s); pagenum(s)
notes(s, "Cómo un peso local se convierte en muchos dólares globales: usar la recaudación como "
          "contrapartida, emitir deuda verde contra ella, y escalar hacia la banca multilateral.")

# =====================================================================
# SLIDE 19 — ÍNDICE DE CIUDAD DESEADA
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "La regla de gasto · 05", color=GREEN)
title(s, [("Índice de Ciudad Deseada (ICD)", SLATE, False)])
dims = [("Espacios verdes",GREEN),("Movilidad",WATER),("Compacidad",SLATE),
        ("Servicios",GOLD),("Equidad",CLAY),("Infraestructura agua",WATER)]
cw, ch, gx, gy = Inches(3.75), Inches(0.8), Inches(0.28), Inches(0.24)
x0, y0 = Inches(0.75), Inches(2.35)
for i,(t,c) in enumerate(dims):
    col,row = i%3, i//3
    x = x0+col*(cw+gx); y = y0+row*(ch+gy)
    rect(s, x, y, cw, ch, fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.12)
    rect(s, x+Inches(0.2), y+Inches(0.2), Inches(0.4), Inches(0.4), fill=c,
         shape=MSO_SHAPE.OVAL)
    text(s, x+Inches(0.2), y+Inches(0.2), Inches(0.4), Inches(0.4),
         [[(str(i+1), 15, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x+Inches(0.8), y, cw-Inches(0.95), ch, [[(t, 17, INK, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
# escala A-E
grades = [("A",GREEN),("B",WATER),("C",GOLD),("D",RGBColor(0xC9,0x7A,0x3A)),("E",CLAY)]
gx0 = Inches(0.75); gy0 = Inches(4.65); gw = Inches(0.85)
for i,(g,c) in enumerate(grades):
    x = gx0 + i*(gw+Inches(0.12))
    rect(s, x, gy0, gw, Inches(0.85), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.14)
    text(s, x, gy0, gw, Inches(0.85), [[(g, 28, WHITE, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, Inches(5.9), Inches(4.6), Inches(6.5), Inches(1.0),
     [[("Escala A (≥90%) … E (<25%).", 18, INK, True, False)],
      [("Menor nota = mayor prioridad: el dinero va a la necesidad.", 17, SLATE, False, False)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=6, line_spacing=1.1)
footer(s); pagenum(s)
notes(s, "La regla de gasto es tan importante como la recaudación. El ICD mide seis dimensiones "
          "y dirige la inversión a los barrios peor puntuados. Es equidad convertida en algoritmo.")

# =====================================================================
# SLIDE (nuevo) — ÍNDICE DE PRIORIZACIÓN DE INFRAESTRUCTURA
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "La regla de gasto en acción · 05", color=GREEN)
title(s, [("Índice de Priorización de Infraestructura", SLATE, False)])
bullets(s, [
    [("Traduce el ICD en un ", 20, INK, False), ("diagnóstico territorial", 20, GREEN, True),
     (": 6 ejes, 10 indicadores, distrito por distrito.", 20, INK, False)],
    [("Cada indicador se compara con un ", 20, INK, False), ("valor deseado", 20, GOLD, True),
     (" y se clasifica en tres niveles.", 20, INK, False)],
    [("Orienta la inversión hacia las ", 20, INK, False), ("brechas reales", 20, CLAY, True),
     (", no hacia la capacidad de gasto.", 20, INK, False)],
], x=Inches(0.75), y=Inches(2.15), w=Inches(7.5), gap=18, size=20)
# escala de 3 niveles
scale = [("ÓPTIMO", GREEN), ("MÍNIMO", GOLD), ("CRÍTICO", CLAY)]
for i,(lab,c) in enumerate(scale):
    chip(s, Inches(8.5)+i*Inches(1.55), Inches(2.2), Inches(1.45), lab, color=c, size=12.5)
# franja de hallazgos
hall = [
    ("Compacidad", "la dimensión más crítica", CLAY),
    ("Movilidad", "la mejor cubierta (80% óptimo)", GREEN),
    ("Equidad", "ningún distrito llega al óptimo", GOLD),
]
hy = Inches(3.05)
for i,(h,d,c) in enumerate(hall):
    x = Inches(8.5)
    rect(s, x, hy, Inches(4.45), Inches(0.92), fill=WHITE, line=CARD_BD, line_w=Pt(1),
         shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
    rect(s, x, hy, Inches(0.13), Inches(0.92), fill=c)
    text(s, x+Inches(0.3), hy, Inches(4.0), Inches(0.92),
         [[(h+": ", 16, c, True, False), (d, 15, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.05)
    hy += Inches(1.05)
# nota inferior
text(s, Inches(0.75), Inches(5.75), Inches(7.5), Inches(0.9),
     [[("Ranking general: ", 16, SLATE, True, False),
       ("Ciudad y La Puntilla lideran; Industrial y Vertientes de Pedemonte muestran las mayores brechas.",
        16, INK, False, False)]], space_after=0, line_spacing=1.15)
text(s, Inches(0.75), Inches(6.45), Inches(7.5), Inches(0.4),
     [[("Ejercicio piloto CIPUV–UTDT · fuente: Municipalidad de Luján de Cuyo e INDEC.", 11, MUTED, False, True)]],
     space_after=0)
footer(s); pagenum(s)
notes(s, "Este es el motor real del ICD: un tablero que puntúa cada distrito en seis ejes contra un "
          "valor deseado. No mide capacidad de pago, mide déficit. Es un ejercicio piloto de validación "
          "metodológica de la CIPUV.")

# =====================================================================
# SLIDE (nuevo) — RANKING DE PRIORIDAD POR DISTRITO (puntaje total)
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "El diagnóstico · 05", color=CLAY)
title(s, [("Ranking de prioridad por distrito", SLATE, False)])
ranking = [
    ("Ciudad", 7.0), ("La Puntilla", 6.0), ("Mayor Drummond", 4.5), ("Vistalba", 4.5),
    ("Chacras de Coria", 4.0), ("Carrodilla", 3.5), ("Agrelo", 3.5), ("Ugarteche", 3.5),
    ("Las Compuertas", 3.0), ("Perdriel", 3.0), ("Potrerillos", 3.0), ("El Carrizal", 2.5),
    ("Cacheuta", 2.0), ("Vertientes de Pedemonte", 1.5), ("Industrial", 1.0),
]
def rk_color(v):
    if v >= 4.0: return GREEN
    if v >= 2.5: return GOLD
    return CLAY
bx0 = Inches(3.15)          # inicio de las barras
bxmax = Inches(8.3)         # ancho máximo (para valor 7)
ry0 = Inches(2.05); rh = Inches(0.27); rgap = Inches(0.045)
for i,(name,v) in enumerate(ranking):
    y = ry0 + i*(rh+rgap)
    c = rk_color(v)
    text(s, Inches(0.75), y, Inches(2.3), rh, [[(name, 12.5, INK, True, False)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    bw = int(bxmax * (v / 7.0))
    rect(s, bx0, y+Emu(15000), Emu(bw), rh-Emu(30000), fill=c,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.35)
    text(s, bx0+Emu(bw)+Inches(0.08), y, Inches(0.9), rh,
         [[(f"{v:g}", 12.5, c, True, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
# leyenda / mensaje a la derecha
rect(s, Inches(11.7), Inches(2.05), Inches(1.4), Inches(5.0), fill=None)
rect(s, Inches(11.65), Inches(2.15), Inches(1.45), Inches(2.15), fill=RGBColor(0xEF,0xF1,0xF4),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
text(s, Inches(11.8), Inches(2.35), Inches(1.2), Inches(2.0),
     [[("Puntaje /11", 12, SLATE, True, False)],
      [("≥4  bien", 12, GREEN, True, False)],
      [("2,5–3,9", 12, GOLD, True, False)],
      [("<2,5  déficit", 12, CLAY, True, False)]], space_after=8, line_spacing=1.1)
text(s, Inches(9.55), Inches(4.7), Inches(3.4), Inches(2.0),
     [[("Menor puntaje = ", 15, INK, False, False), ("mayor prioridad", 15, CLAY, True, False),
       (" de inversión.", 15, INK, False, False)],
      [("", 6, INK, False, False)],
      [("El índice manda el dinero a Industrial, Vertientes y Cacheuta —no a Ciudad.", 14, SLATE, False, True)]],
     space_after=8, line_spacing=1.15)
footer(s); pagenum(s)
notes(s, "El resultado en una imagen: Ciudad y La Puntilla, ya consolidadas, puntúan alto; los distritos "
          "del oeste y el periurbano (Industrial, Vertientes, Cacheuta) puntúan bajo y son, por diseño, "
          "los primeros destinatarios de la inversión. Puntaje sobre 11 puntos.")

# =====================================================================
# SLIDE (nuevo) — TABLA SÍNTESIS (11 subíndices, tabla real CIPUV)
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Síntesis · 05", color=GREEN)
title(s, [("Tabla Síntesis: los 11 subíndices por distrito", SLATE, False)], size=30)
# leyenda de 3 niveles (arriba a la derecha)
for i,(lab,c) in enumerate([("ÓPTIMO", GREEN), ("MÍNIMO", GOLD), ("CRÍTICO", CLAY)]):
    chip(s, Inches(8.9)+i*Inches(1.45), Inches(1.12), Inches(1.35), lab, color=c, size=11.5)
# imagen de la tabla (1318x600) centrada, ajustada sobre el pie
_tw = Inches(11.6); _th = Emu(int(_tw * 600 / 1318))
s.shapes.add_picture(os.path.join(ASSETS, "tabla-sintesis-infraestructura.png"),
                     Inches((13.333-11.6)/2), Inches(1.7), width=_tw, height=_th)
text(s, Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.32),
     [[("Cada indicador se pondera sobre 11 puntos: Óptimo 1 · Mínimo 0,5 · Crítico 0 · s/d computa 0.  ", 11.5, MUTED, False, True),
       ("Fuente: elaboración propia CIPUV–UTDT.", 11.5, SLATE2, False, True)]], space_after=0)
pagenum(s)
notes(s, "La tabla completa que sostiene el ranking: 11 subíndices por distrito, cada celda clasificada "
          "en Óptimo (verde), Mínimo (amarillo) o Crítico (rojo). El indicador 1.1 (infraestructura verde "
          "pública) queda pendiente por falta de datos. Es la tabla original del ejercicio piloto CIPUV.")

# =====================================================================
# SLIDE 20 — PRUEBA DE EQUIDAD
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "La prueba de equidad · 05", color=CLAY)
title(s, [("¿Dónde se necesita vs. dónde se invierte?", SLATE, False)])
bullets(s, [
    [("Las 2 mayores inversiones caen en distritos que ", 21, INK, False), ("ya puntúan bien", 21, WATER, True), (" en verde.", 21, INK, False)],
    [("Los más necesitados ", 21, INK, False), ("(Vistalba, Las Compuertas, Industrial)", 21, CLAY, True), (" no reciben obra verde.", 21, INK, False)],
    [("La palanca: ", 21, INK, False), ("reasignar las microplazas (ARS 1.200 M) a distritos D/E", 21, GREEN, True), (".", 21, INK, False)],
], w=Inches(11.6), gap=22)
text(s, Inches(0.75), Inches(5.9), Inches(11.7), Inches(0.7),
     [[("El ICD no describe: corrige. Alinea la inversión con la necesidad, no con la capacidad.",
        18, GOLD, True, True)]], space_after=0, line_spacing=1.1)
footer(s); pagenum(s)
notes(s, "La honestidad del caso: sin regla, el dinero va a donde ya está bien. El ICD existe "
          "precisamente para forzar la corrección hacia los distritos D y E.")

# =====================================================================
# SLIDE 21 — LOS NÚMEROS REALES
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Los números reales · 06", color=GOLD)
title(s, [("Potencial de recaudación", SLATE, False)])
nums = [
    ("~USD 1 M", "YA se recauda hoy (2 instrumentos operativos)", GREEN),
    ("~USD 3 M", "potencial del paquete completo (base)", GOLD),
    ("USD 56 M", "brecha de espacio verde → la cierra el bono verde", WATER),
]
cw = Inches(3.85); gx = Inches(0.28); x0 = Inches(0.75); y0 = Inches(2.7)
for i,(big,small,c) in enumerate(nums):
    x = x0 + i*(cw+gx)
    rect(s, x, y0, cw, Inches(2.9), fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
    rect(s, x, y0, cw, Pt(6), fill=c)
    text(s, x+Inches(0.25), y0+Inches(0.55), cw-Inches(0.5), Inches(1.0),
         [[(big, 40, c, True, False)]], align=PP_ALIGN.CENTER, space_after=0)
    text(s, x+Inches(0.3), y0+Inches(1.7), cw-Inches(0.6), Inches(1.0),
         [[(small, 17, SLATE, False, False)]], align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.15)
text(s, Inches(0.75), Inches(6.15), Inches(11.7), Inches(0.5),
     [[("De ~1 a ~3 con el paquete; y del recurso propio, un salto a decenas con el bono verde.",
        17, MUTED, False, True)]], align=PP_ALIGN.CENTER, space_after=0)
footer(s); pagenum(s)
notes(s, "La escala en tres cifras: lo que ya hay, el potencial propio y la brecha que sólo cierra "
          "el apalancamiento. El fideicomiso es el puente entre la segunda y la tercera.")

# =====================================================================
# SLIDE 22 — QUIÉN PAGA / QUIÉN SE BENEFICIA
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Equidad · 06", color=CLAY)
title(s, [("¿Quién paga y quién se beneficia?", SLATE, False)])
rect(s, Inches(0.75), Inches(2.4), Inches(5.7), Inches(2.3), fill=CLAY_BG, shadow=True,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
text(s, Inches(1.05), Inches(2.65), Inches(5.1), Inches(0.5), [[("PAGA", 18, CLAY, True, False)]], space_after=0)
text(s, Inches(1.05), Inches(3.2), Inches(5.1), Inches(1.3),
     [[("Quien captura la plusvalía: desarrolladores y dueños de suelo ocioso (bases progresivas).", 19, INK, False, False)]],
     space_after=0, line_spacing=1.15)
rect(s, Inches(6.85), Inches(2.4), Inches(5.7), Inches(2.3), fill=GREEN_BG, shadow=True,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
text(s, Inches(7.15), Inches(2.65), Inches(5.1), Inches(0.5), [[("SE BENEFICIA", 18, GREEN, True, False)]], space_after=0)
text(s, Inches(7.15), Inches(3.2), Inches(5.1), Inches(1.3),
     [[("El distrito con menor acceso a verde, agua y servicios (vía ICD).", 19, INK, False, False)]],
     space_after=0, line_spacing=1.15)
rect(s, Inches(0.75), Inches(5.05), Inches(11.8), Inches(1.35), fill=SLATE_DK,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.07)
text(s, Inches(1.1), Inches(5.05), Inches(11.1), Inches(1.35),
     [[("Con suelo inelástico (límite hídrico), la carga recae en el terrateniente → ", 19, WHITE, False, False),
       ("un caso deliberadamente progresivo.", 19, GOLD_LT, True, False)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.12)
footer(s); pagenum(s)
notes(s, "La respuesta a la pregunta de examen: como el suelo es inelástico por el límite del agua, "
          "el impuesto no se traslada, recae en el dueño de la tierra. Es progresivo.")

# =====================================================================
# SLIDE 23 — CO-BENEFICIOS
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "No sólo adaptación · 06", color=GREEN)
title(s, [("Co-beneficios ", SLATE, False), ("low-carbon", GREEN, False)])
bullets(s, [
    [("Premian el desarrollo compacto y orientado al transporte (TOD): ", 21, INK, False),
     ("el metrotranvía concentra edificabilidad.", 21, GREEN, True)],
    [("Menos viajes en auto: ", 21, INK, False), ("menor consumo energético e hídrico per cápita.", 21, WATER, True)],
    [("Se evita convertir suelo agrícola y de pedemonte: ", 21, INK, False), ("carbono + recarga preservados.", 21, GOLD, True)],
], w=Inches(11.6), gap=22)
text(s, Inches(0.75), Inches(5.9), Inches(11.7), Inches(0.7),
     [[("Adaptación y mitigación: la misma palanca, accionada una sola vez.", 19, GOLD, True, True)]],
     space_after=0, line_spacing=1.1)
footer(s); pagenum(s)
notes(s, "El bonus: los mismos instrumentos que financian la adaptación también reducen emisiones. "
          "No hay que elegir entre adaptar y mitigar; aquí es la misma decisión.")

# =====================================================================
# SLIDE 24 — PARA DISCUTIR
# =====================================================================
s = slide(); base(s, bg=SLATE_DK)
rect(s, 0, 0, Inches(0.32), EMU_H, fill=GOLD)
text(s, Inches(0.95), Inches(0.85), Inches(11), Inches(0.5),
     [[("PARA DISCUTIR", 16, GOLD_LT, True, False)]], space_after=0)
text(s, Inches(0.95), Inches(1.4), Inches(11.5), Inches(0.9),
     [[("Cinco preguntas para el debate", 34, WHITE, True, False)]], space_after=0)
qs = [
    "¿La captura de valor es regresiva o progresiva? ¿De qué depende?",
    "¿Por qué un fideicomiso y no el presupuesto general?",
    "¿Qué pasa si el take-up de derechos de construcción es bajo?",
    "¿Lo replicarías en tu ciudad? ¿Qué instrumento primero?",
    "¿Cómo se sostiene políticamente entre gestiones?",
]
y = Inches(2.6)
for i,q in enumerate(qs):
    rect(s, Inches(0.95), y, Inches(0.55), Inches(0.55), fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, Inches(0.95), y, Inches(0.55), Inches(0.55), [[("?", 24, WHITE, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, Inches(1.75), y, Inches(10.6), Inches(0.6),
         [[(q, 20, RGBColor(0xE4,0xE7,0xEC), False, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.05)
    y += Inches(0.82)
notes(s, "Preguntas para abrir el debate. La primera y la última suelen dar la mejor discusión: "
          "incidencia del impuesto y sostenibilidad política.")

# =====================================================================
# SLIDE 25 — PARA LLEVARSE
# =====================================================================
s = slide(); base(s); brandbar(s); kicker(s, "Para llevarse")
title(s, [("Cuatro ideas", SLATE, False)])
takes = [
    ("Una máquina replicable", "capturar → blindar → apalancar → redistribuir", WATER),
    ("Ya funciona", "~USD 1 M/año se recauda; el paquete llega a ~USD 3 M", GREEN),
    ("El fideicomiso es la llave", "convierte recurso local en contrapartida y repago", GOLD),
    ("Equidad por diseño", "el índice manda el dinero a donde más se necesita", CLAY),
]
cw2, ch2, gx2, gy2 = Inches(5.85), Inches(1.35), Inches(0.28), Inches(0.26)
x0, y0 = Inches(0.75), Inches(2.25)
for i,(h,d,c) in enumerate(takes):
    col,row = i%2, i//2
    x = x0+col*(cw2+gx2); yy = y0+row*(ch2+gy2)
    rect(s, x, yy, cw2, ch2, fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
    rect(s, x, yy, Inches(0.14), ch2, fill=c)
    text(s, x+Inches(0.35), yy+Inches(0.18), cw2-Inches(0.6), Inches(0.5),
         [[(f"{i+1}.  ", 20, c, True, False),(h, 20, INK, True, False)]], space_after=0)
    text(s, x+Inches(0.35), yy+Inches(0.72), cw2-Inches(0.6), Inches(0.55),
         [[(d, 16, SLATE, False, False)]], space_after=0, line_spacing=1.08)
# cita
rect(s, Inches(0.75), Inches(5.5), Inches(11.8), Inches(1.0), fill=SLATE_DK,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
text(s, Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.0),
     [[("“Una transición justa no se mide por el tamaño del cheque, sino por dónde aterriza.”",
        20, WHITE, True, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.1)
footer(s); pagenum(s)
notes(s, "Cierre. Cuatro ideas para llevarse y una frase: lo que define una transición justa no es "
          "cuánto dinero, sino dónde cae. Gracias.")

# ---------------- guardar ----------------
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "lujan-de-cuyo-gestion-suelo.pptx")
prs.save(out)
print("OK ->", out, "· slides:", len(prs.slides._sldIdLst))
