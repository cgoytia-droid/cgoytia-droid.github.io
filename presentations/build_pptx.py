#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Genera "Ciudades Habitables y Resilientes" como PowerPoint editable (.pptx).
Basado en: World Bank & GFDRR (2025), Handbook for Livable and Resilient Cities.
Misma narrativa, identidad visual y notas del orador que la versión interactiva HTML.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- paleta (sitio de la autora) ----------
SLATE   = RGBColor(0x4E,0x58,0x69)
SLATE2  = RGBColor(0x6F,0x7D,0x94)
SLATE_DK= RGBColor(0x3A,0x42,0x50)
GOLD    = RGBColor(0xB8,0x86,0x0B)
GOLD_LT = RGBColor(0xD4,0xC0,0x88)
INK     = RGBColor(0x2B,0x30,0x38)
PAPER   = RGBColor(0xF7,0xF8,0xFA)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
MUTED   = RGBColor(0x7A,0x82,0x8E)
ENV     = RGBColor(0x5A,0x8F,0x6B)
PEOPLE  = RGBColor(0xA9,0x89,0x7A)
INFRA   = RGBColor(0xB1,0x45,0x6A)
ECON    = RGBColor(0x2F,0x8F,0x9D)
CARD_BD = RGBColor(0xE6,0xE8,0xEC)
FROM_BG = RGBColor(0xF1,0xE9,0xE6)
TO_BG   = RGBColor(0xE7,0xF0,0xEA)

FONT = "Verdana"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def _set_font(run, size, color, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = name; run.font.color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, round_=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    if round_ is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = round_
        except Exception: pass
    return sp

def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
        {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val':'28323E'})
    alpha = clr.makeelement(qn('a:alpha'), {'val':'22000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            txt, size, color, bold, italic = (list(seg)+[False,False])[:5]
            r = p.add_run(); r.text = txt; _set_font(r, size, color, bold, italic)
    return tb

def kicker(s, label, x=Inches(0.7), y=Inches(0.55), color=GOLD):
    bar = rect(s, x, y+Emu(20000), Pt(4), Inches(0.28), fill=color)
    text(s, x+Inches(0.12), y, Inches(9), Inches(0.34),
         [[(label.upper(), 12, color, True, False)]], anchor=MSO_ANCHOR.MIDDLE)

def title(s, parts, x=Inches(0.7), y=Inches(0.95), w=Inches(12), size=34):
    # parts: list of (txt, color, italic)
    runs=[[(t, size, c, True, it) for (t,c,it) in parts]]
    text(s, x, y, w, Inches(1.2), runs, line_spacing=1.0, space_after=0)

def brandbar(s):
    rect(s, 0, 0, EMU_W, Pt(6), fill=GOLD)

def footer(s):
    text(s, Inches(0.7), Inches(7.12), Inches(9), Inches(0.3),
         [[("Handbook for Livable & Resilient Cities  ·  ", 9, MUTED, False, False),
           ("World Bank / GFDRR 2025", 9, SLATE2, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)

def source(s, txt, y=Inches(6.55)):
    text(s, Inches(0.7), y, Inches(11.9), Inches(0.5),
         [[(txt, 10, MUTED, False, True)]], space_after=0, line_spacing=1.05)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def base(s, dark=False):
    bg = s.background; bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_DK if dark else PAPER
    brandbar(s)

def icon_circle(s, cx, cy, d, color, glyph):
    rect(s, cx-d/2, cy-d/2, d, d, fill=color, shape=MSO_SHAPE.OVAL)
    text(s, cx-d/2, cy-d/2, d, d, [[(glyph, 22, WHITE, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

# ============================================================
# 1. PORTADA
# ============================================================
s = slide(); base(s, dark=True)
# fondo: bandas tipo "delta" con rectángulos translúcidos
rect(s, 0, 0, EMU_W, EMU_H, fill=SLATE_DK)
for i,(yy,hh,col) in enumerate([(2.4,0.5,SLATE2),(3.1,0.32,SLATE2),(4.0,0.22,SLATE2)]):
    b = rect(s, Inches(-0.5), Inches(yy), Inches(14), Inches(hh), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
    b.rotation = -6
rect(s, Inches(6.3), Inches(-1.2), Inches(8), Inches(8), fill=SLATE2, shape=MSO_SHAPE.OVAL)._element  # halo
# tarjeta
card = rect(s, Inches(0.9), Inches(1.7), Inches(8.6), Inches(3.9), fill=WHITE, shadow=True,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.03)
rect(s, Inches(0.9), Inches(1.7), Pt(6), Inches(3.9), fill=GOLD)
text(s, Inches(1.25), Inches(2.0), Inches(8), Inches(0.35),
     [[("PLANIFICACIÓN URBANA INFORMADA POR EL RIESGO", 12, GOLD, True, False)]])
text(s, Inches(1.25), Inches(2.45), Inches(8), Inches(1.7),
     [[("Ciudades Habitables", 46, SLATE, True, False)],
      [("y Resilientes", 46, SLATE, True, False)]], line_spacing=0.95, space_after=0)
text(s, Inches(1.25), Inches(4.35), Inches(8), Inches(0.7),
     [[("Integrar la información de ", 16, INK, False, False),
       ("amenazas y riesgos", 16, SLATE_DK, True, False),
       (" en la planificación urbana", 16, INK, False, False)]], line_spacing=1.1)
text(s, Inches(1.25), Inches(4.95), Inches(8), Inches(0.7),
     [[("Prof. Cynthia Goytia (PhD)", 13, SLATE_DK, True, False),
       ("  ·  Economía Urbana · UTDT — CIPUV", 13, SLATE, False, False)],
      [("A partir del Handbook for Livable and Resilient Cities · World Bank · GFDRR · 2025",
        11, MUTED, False, True)]], line_spacing=1.2)
notes(s, "Bienvenida. Arranco con calidez y una afirmación grande: 'Vivimos el siglo de las "
        "ciudades. Dónde y cómo se construyan en las próximas tres décadas va a decidir buena "
        "parte del futuro del clima, de la equidad y de la prosperidad.'\n"
        "Encuadre: esta charla se apoya en el nuevo manual del Banco Mundial y GFDRR (2025). No "
        "es un compendio técnico: es una invitación a cambiar la forma en que planificamos.\n"
        "Tiempo total objetivo: 20–25 min. Esta slide: ~45 s. Transición: 'Empecemos por entender "
        "por qué las ciudades importan tanto.'")

# ============================================================
# 2. GANCHO: EL SIGLO DE LAS CIUDADES
# ============================================================
s = slide(); base(s)
kicker(s, "El siglo de las ciudades")
title(s, [("Las ciudades concentran el futuro", SLATE, False)])
stats = [("70%","de la población mundial vivirá en ciudades hacia 2050 (hoy 56%) · ~7.000 millones", GOLD),
         ("80%","del PBI global se produce en áreas urbanas", GOLD),
         ("70%","de las emisiones de GEI ocurren en ciudades", GOLD),
         ("70%","del consumo energético mundial es urbano", GOLD)]
x0,y0,w,h,gx,gy = Inches(0.7),Inches(2.0),Inches(5.85),Inches(1.85),Inches(0.3),Inches(0.3)
for i,(n,l,col) in enumerate(stats):
    cx = x0 + (i%2)*(w+gx); cy = y0 + (i//2)*(h+gy)
    rect(s, cx, cy, w, h, fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.04)
    rect(s, cx, cy, w, Pt(4), fill=col)
    text(s, cx+Inches(0.3), cy+Inches(0.18), Inches(3), Inches(0.8),
         [[(n, 40, SLATE, True, False)]], space_after=0)
    text(s, cx+Inches(0.3), cy+Inches(1.0), w-Inches(0.6), Inches(0.75),
         [[(l, 12.5, RGBColor(0x5C,0x63,0x6E), False, False)]], line_spacing=1.1)
source(s, "Fuente: World Bank 2022a · Our World in Data 2025 (Fig. 1.1).")
footer(s)
notes(s, "Mensaje: la ciudad es el motor de la civilización — riqueza, innovación, empleo, "
        "cultura. Recalcar el salto 56% → 70%.\nPausa después de '80% del PBI': 'Si las ciudades "
        "fueran un país, serían LA economía mundial.'\nCierre del beat: 'Pero esa misma "
        "concentración tiene una contracara.' — puente a la próxima slide. ~1:30.")

# ============================================================
# 3. LA CONTRACARA: CONCENTRAN EL RIESGO
# ============================================================
s = slide(); base(s)
kicker(s, "La contracara")
title(s, [("…y también concentran el ", SLATE, False), ("riesgo", GOLD, False)])
# pill
pill = rect(s, Inches(0.7), Inches(2.05), Inches(3.6), Inches(0.45), fill=SLATE,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
text(s, Inches(0.7), Inches(2.05), Inches(3.6), Inches(0.45),
     [[("679 GRANDES CIUDADES", 12, WHITE, True, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, Inches(0.7), Inches(2.65), Inches(6.6), Inches(0.6),
     [[("(≥500.000 hab.) analizadas frente a ciclones, inundaciones, sequías, sismos, "
        "deslizamientos y erupciones.", 14, INK, False, False)]], line_spacing=1.2)
bullets = [("59%"," está expuesta a al menos una amenaza"),
           ("189 ciudades"," enfrentan dos o más"),
           ("26 ciudades"," enfrentan tres o más")]
yy = Inches(3.5)
for b,rest in bullets:
    rect(s, Inches(0.85), yy+Inches(0.1), Inches(0.12), Inches(0.12), fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.15), yy, Inches(6), Inches(0.4),
         [[(b, 16, SLATE_DK, True, False), (rest, 16, INK, False, False)]], space_after=0)
    yy = yy + Inches(0.5)
text(s, Inches(0.7), Inches(5.35), Inches(6.6), Inches(0.9),
     [[("El cambio climático golpea de forma desproporcionada a los más pobres, mujeres, "
        "niñas y niños, personas mayores y minorías.", 12.5, MUTED, False, False)]], line_spacing=1.25)
# big stat right
rect(s, Inches(8.0), Inches(2.4), Inches(4.5), Inches(2.6), fill=WHITE, line=CARD_BD, line_w=Pt(1),
     shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.04)
rect(s, Inches(8.0), Inches(2.4), Inches(4.5), Pt(5), fill=INFRA)
text(s, Inches(8.0), Inches(2.75), Inches(4.5), Inches(1.3),
     [[("59%", 80, INFRA, True, False)]], align=PP_ALIGN.CENTER, space_after=0)
text(s, Inches(8.25), Inches(4.05), Inches(4.0), Inches(0.9),
     [[("de las grandes ciudades del mundo, expuestas a una o más amenazas naturales",
        13, RGBColor(0x5C,0x63,0x6E), False, False)]], align=PP_ALIGN.CENTER, line_spacing=1.15)
source(s, "Fuente: United Nations 2018 (Fig. 1.2); World Bank 2023b.")
footer(s)
notes(s, "Tono: bajar la energía, hacerlo humano. 'Detrás de cada porcentaje hay personas.' "
        "Subrayar la inequidad: el riesgo no se reparte parejo — lo paga primero quien menos "
        "tiene.\nConectar con la región: en América Latina gran parte del crecimiento es informal "
        "y en suelo expuesto (agregar ejemplo del AMBA si el público es local).\nTransición: '¿Y "
        "por qué crece tanto el riesgo? Acá viene lo incómodo.' ~1:30.")

# ============================================================
# 4. LA PARADOJA
# ============================================================
s = slide(); base(s)
kicker(s, "La paradoja")
title(s, [("La urbanización ", SLATE, False), ("misma", GOLD, True), (" produce riesgo", SLATE, False)])
para = [("La expansión urbana supera el crecimiento de la población hasta en un 50%: "
         "transforma el suelo, fragmenta ecosistemas y concentra gente y activos en zonas peligrosas."),
        ("El desarrollo urbano genera nuevos riesgos: inundación por superficies impermeables, "
         "deslizamientos en laderas desestabilizadas."),
        ("Los desastres climáticos se duplicaron con creces en los últimos 40 años.")]
yy = Inches(2.05)
for p in para:
    rect(s, Inches(0.85), yy+Inches(0.12), Inches(0.14), Inches(0.14), fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.2), yy, Inches(11.2), Inches(0.8),
         [[(p, 16, INK, False, False)]], line_spacing=1.2, space_after=0)
    yy = yy + Inches(0.95)
# alerta ×20
box = rect(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.15), fill=FROM_BG,
           line=RGBColor(0xE3,0xCD,0xC6), line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
text(s, Inches(1.0), Inches(5.05), Inches(9.5), Inches(1.15),
     [[("⚠  El daño por inundaciones fluviales podría multiplicarse hacia fin de siglo si no se actúa.",
        15, SLATE_DK, True, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
text(s, Inches(10.4), Inches(5.05), Inches(2.0), Inches(1.15),
     [[("×20", 44, GOLD, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
source(s, "Fuente: CRED & UNDRR 2019; Winsemius et al. 2016.", y=Inches(6.35))
notes(s, "Idea central y contraintuitiva: la ciudad no es un escenario pasivo donde 'ocurre' el "
        "desastre. La forma en que urbanizamos CREA o REDUCE el riesgo. Es una decisión.\nDejar "
        "respirar el ×20 — es el dato que la gente se lleva. 'No es destino. Es la suma de "
        "decisiones de planificación que todavía estamos a tiempo de cambiar.'\nTransición: 'Si la "
        "urbanización crea el riesgo, el problema de fondo está en cómo planificamos.' ~1:20.")

# ============================================================
# 5. EL PROBLEMA DE FONDO (oscura)
# ============================================================
s = slide(); base(s, dark=True)
kicker(s, "El problema de fondo", color=GOLD_LT)
title(s, [("Planificamos ", WHITE, False), ("sin mirar el riesgo", GOLD_LT, True)])
text(s, Inches(0.7), Inches(2.4), Inches(8.8), Inches(1.0),
     [[("La planificación urbana convencional suele ignorar las amenazas existentes y las "
        "vulnerabilidades de largo plazo.", 20, RGBColor(0xDF,0xE3,0xE9), False, False)]], line_spacing=1.25)
text(s, Inches(0.7), Inches(3.7), Inches(9.2), Inches(1.0),
     [[("Resultado: una ", 20, RGBColor(0xDF,0xE3,0xE9), False, False),
       ("acumulación silenciosa de riesgo", 20, GOLD_LT, True, False),
       (" que deja a las ciudades mal preparadas para los shocks presentes y futuros.",
        20, RGBColor(0xDF,0xE3,0xE9), False, False)]], line_spacing=1.25)
text(s, Inches(0.7), Inches(5.2), Inches(9), Inches(0.7),
     [[("Hace falta un enfoque nuevo.", 26, WHITE, True, False)]], space_after=0)
notes(s, "Slide de tensión máxima — fondo oscuro a propósito. Hablar más lento. El riesgo se "
        "acumula como una deuda que no vemos en el balance, hasta que llega el evento y la "
        "pagamos toda junta.\nPlantar la pregunta que ordena la charla: '¿Y si la ciudad, en "
        "lugar de ser el problema, fuera la solución?' — apenas insinuarla. ~1:00.")

# ============================================================
# 6. CAMBIO DE PARADIGMA (ES.1)
# ============================================================
s = slide(); base(s)
kicker(s, "El cambio de paradigma · Figura ES.1")
title(s, [("De motor del riesgo a ", SLATE, False), ("catalizador de resiliencia", GOLD, False)], size=30)
# from / arrow / to
fw = Inches(5.2)
rect(s, Inches(0.7), Inches(1.95), fw, Inches(1.25), fill=FROM_BG, line=RGBColor(0xE3,0xCD,0xC6),
     line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
text(s, Inches(0.95), Inches(2.1), fw-Inches(0.5), Inches(1.05),
     [[("Ciudades como focos de riesgo", 15, INFRA, True, False)],
      [("La urbanización como motor de exposición y riesgo", 13, INK, False, False)]], line_spacing=1.15)
arrow = rect(s, Inches(6.05), Inches(2.25), Inches(1.0), Inches(0.65), fill=GOLD, shape=MSO_SHAPE.RIGHT_ARROW)
rect(s, Inches(7.15), Inches(1.95), fw, Inches(1.25), fill=TO_BG, line=RGBColor(0xC5,0xDD,0xCD),
     line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
text(s, Inches(7.4), Inches(2.1), fw-Inches(0.5), Inches(1.05),
     [[("Ciudades como parte de la solución", 15, ENV, True, False)],
      [("La urbanización como motor de capacidad y resiliencia", 13, INK, False, False)]], line_spacing=1.15)
text(s, Inches(0.7), Inches(3.4), Inches(11), Inches(0.3),
     [[("Cinco principios que sostienen el giro:", 12.5, MUTED, False, False)]], space_after=0)
principles = [
    ("01","Las amenazas urbanas son multifacéticas, no sólo naturales."),
    ("02","Planificar exige integrar riesgos presentes y futuros."),
    ("03","La planificación sistémica e informada por el riesgo es la base."),
    ("04","Integra planificación + DRM + gestión climática."),
    ("05","Se puede avanzar aun con datos y capacidades limitadas.")]
pw = Inches(2.28); px = Inches(0.7); gap = Inches(0.12)
for i,(n,t) in enumerate(principles):
    cx = px + i*(pw+gap)
    rect(s, cx, Inches(3.8), pw, Inches(2.0), fill=WHITE, line=CARD_BD, line_w=Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
    rect(s, cx, Inches(5.72), pw, Pt(4), fill=SLATE2)
    text(s, cx+Inches(0.2), Inches(3.95), pw-Inches(0.4), Inches(0.5),
         [[(n, 22, GOLD, True, False)]], space_after=0)
    text(s, cx+Inches(0.2), Inches(4.5), pw-Inches(0.4), Inches(1.2),
         [[(t, 11.5, INK, False, False)]], line_spacing=1.15, space_after=0)
source(s, "Fuente: World Bank / GFDRR 2025, Figura ES.1.")
footer(s)
notes(s, "Columna vertebral de la charla. El giro: la urbanización no es el enemigo; mal "
        "gestionada es riesgo, bien gestionada es resiliencia.\nRecorrer los 5 principios con "
        "ritmo. Detenerse en el 5: el más esperanzador y político — 'no hay excusa de no tenemos "
        "datos: se empieza con lo que hay'. Lo retomamos.\nTransición: '¿Hacia qué ciudad queremos "
        "ir?' ~1:40.")

# ============================================================
# 7. DEFINICIÓN
# ============================================================
s = slide(); base(s)
kicker(s, "¿Qué buscamos?")
title(s, [("Una ciudad habitable y resiliente", SLATE, False)])
text(s, Inches(0.7), Inches(2.3), Inches(11.5), Inches(1.4),
     [[("Es aquella —y su entorno— donde se promueven el ", 22, INK, False, False),
       ("crecimiento urbano verde", 22, ENV, True, False),
       (", la ", 22, INK, False, False),
       ("inclusión social", 22, PEOPLE, True, False),
       (", entornos construidos ", 22, INK, False, False),
       ("resilientes", 22, INFRA, True, False),
       (" y la ", 22, INK, False, False),
       ("prosperidad compartida", 22, ECON, True, False),
       (".", 22, INK, False, False)]], line_spacing=1.3)
text(s, Inches(0.7), Inches(4.6), Inches(11), Inches(1.2),
     [[("Donde el crecimiento planificado y sostenible asegura acceso a entornos saludables, "
        "vivienda asequible, servicios básicos, empleo y transporte de bajo carbono.",
        16, MUTED, False, False)]], line_spacing=1.3)
footer(s)
notes(s, "Slide-respiro, positiva. Leer la definición como una visión, no como un glosario. "
        "Recalcar 'vivienda asequible' y 'prosperidad compartida' — el puente entre resiliencia y "
        "equidad, el corazón de la agenda.\nTransición: 'Esta visión se ordena en cuatro "
        "objetivos.' ~0:50.")

# ============================================================
# 8. LOS 4 OBJETIVOS (ES.2)
# ============================================================
s = slide(); base(s)
kicker(s, "Cuatro objetivos · Figura ES.2")
title(s, [("Metas de la planificación informada por el riesgo", SLATE, False)], size=30)
goals = [("AMBIENTE","Crecimiento urbano verde", ENV, "🌱"),
         ("PERSONAS","Inclusión social", PEOPLE, "👥"),
         ("INFRAESTRUCTURA","Entorno construido resiliente", INFRA, "🏙"),
         ("ECONOMÍA","Prosperidad compartida", ECON, "📈")]
gw = Inches(2.85); gx = Inches(0.7); gap = Inches(0.23)
for i,(cat,gt,col,gl) in enumerate(goals):
    cx = gx + i*(gw+gap)
    rect(s, cx, Inches(2.2), gw, Inches(3.1), fill=WHITE, line=CARD_BD, line_w=Pt(1),
         shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
    icon_circle(s, cx+gw/2, Inches(3.0), Inches(1.05), col, gl)
    text(s, cx+Inches(0.15), Inches(3.75), gw-Inches(0.3), Inches(0.3),
         [[(cat, 11, MUTED, True, False)]], align=PP_ALIGN.CENTER, space_after=0)
    text(s, cx+Inches(0.2), Inches(4.05), gw-Inches(0.4), Inches(1.1),
         [[(gt, 17, col, True, False)]], align=PP_ALIGN.CENTER, line_spacing=1.1, space_after=0)
source(s, "Fuente: World Bank / GFDRR 2025, Figura ES.2. Requiere colaboración entre planificación "
          "urbana, gestión del riesgo de desastres (DRM) y gestión del riesgo climático (CRM).")
footer(s)
notes(s, "Las cuatro patas de la mesa: Ambiente, Personas, Infraestructura, Economía. Ninguna se "
        "sostiene sola; la resiliencia aparece cuando se integran.\nClave de gobernanza: NO lo "
        "hace sólo el planificador urbano. Requiere sentar en la misma mesa a planificación + DRM "
        "+ CRM. Romper los silos es la mitad del trabajo.\nTransición: '¿Cómo se traduce en la "
        "práctica? Primero, un cambio de actitud.' ~1:10.")

# ============================================================
# 9. REACTIVO -> PROACTIVO
# ============================================================
s = slide(); base(s)
kicker(s, "La solución, en una frase")
title(s, [("De ", SLATE, False), ("reaccionar", INFRA, True), (" a ", SLATE, False), ("anticipar", GOLD, True)])
fw = Inches(5.2)
rect(s, Inches(0.7), Inches(2.1), fw, Inches(1.4), fill=FROM_BG, line=RGBColor(0xE3,0xCD,0xC6),
     line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
text(s, Inches(0.95), Inches(2.3), fw-Inches(0.5), Inches(1.1),
     [[("Enfoque reactivo", 16, INFRA, True, False)],
      [("Respuesta a la emergencia y recuperación después del desastre", 14, INK, False, False)]], line_spacing=1.2)
rect(s, Inches(6.05), Inches(2.5), Inches(1.0), Inches(0.65), fill=GOLD, shape=MSO_SHAPE.RIGHT_ARROW)
rect(s, Inches(7.15), Inches(2.1), fw, Inches(1.4), fill=TO_BG, line=RGBColor(0xC5,0xDD,0xCD),
     line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
text(s, Inches(7.4), Inches(2.3), fw-Inches(0.5), Inches(1.1),
     [[("Enfoque proactivo", 16, ENV, True, False)],
      [("Prevenir y reducir los impactos antes de que ocurran", 14, INK, False, False)]], line_spacing=1.2)
text(s, Inches(0.7), Inches(4.1), Inches(11.6), Inches(1.4),
     [[("Incorporar la prevención al desarrollo urbano permite anticipar amenazas, minimizar "
        "pérdidas humanas y económicas y fortalecer la capacidad adaptativa de personas, "
        "ecosistemas, infraestructura y medios de vida.", 17, INK, False, False)]], line_spacing=1.3)
footer(s)
notes(s, "Mensaje económico (tu fuerte): cada peso en prevención rinde mucho más que el gasto en "
        "reconstrucción. Reaccionar es caro y tarde; anticipar es barato y a tiempo.\nTransición: "
        "'Para anticipar, el manual propone un instrumento simple y potente.' ~0:55.")

# ============================================================
# 10. MARCO RCP
# ============================================================
s = slide(); base(s)
kicker(s, "El instrumento clave")
title(s, [("El marco ", SLATE, False), ("RCP", GOLD, False),
          (": dónde y cómo construir", SLATE, False)], size=30)
rcp = [("R","Restringir","Evitar el desarrollo en zonas de alto peligro: áreas de no construcción, "
        "control de la expansión en suelo expuesto.", INFRA),
       ("C","Condicionar","Permitir construir, pero bajo reglas: retiros, códigos de edificación, "
        "estándares de diseño resiliente.", GOLD),
       ("P","Promover","Incentivar el desarrollo seguro y la transformación resiliente donde el "
        "riesgo es manejable.", ECON)]
cw = Inches(3.85); cx0 = Inches(0.7); gap = Inches(0.18)
for i,(L,ti,p,col) in enumerate(rcp):
    cx = cx0 + i*(cw+gap)
    rect(s, cx, Inches(2.2), cw, Inches(3.2), fill=col, shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
    text(s, cx+Inches(0.3), Inches(2.4), Inches(2), Inches(1.0),
         [[(L, 54, WHITE, True, False)]], space_after=0)
    text(s, cx+Inches(0.3), Inches(3.55), cw-Inches(0.6), Inches(0.5),
         [[(ti, 20, WHITE, True, False)]], space_after=0)
    text(s, cx+Inches(0.3), Inches(4.1), cw-Inches(0.6), Inches(1.2),
         [[(p, 13, WHITE, False, False)]], line_spacing=1.2, space_after=0)
source(s, "El marco Restringir–Condicionar–Promover (RCP) es la columna operativa del Handbook "
          "(Partes I y III): traduce la información de riesgo en decisiones de planificación.")
footer(s)
notes(s, "El RCP es lo más 'llevable a casa' de toda la charla. Tres verbos que cualquier "
        "intendente entiende: dónde NO construir (Restringir), dónde construir CON cuidado "
        "(Condicionar), dónde EMPUJAR desarrollo resiliente (Promover).\nEjemplo: llanura de "
        "inundación → Restringir; ladera con buen drenaje y código → Condicionar; área central "
        "segura y con servicios → Promover densidad.\nTransición: 'Lo mejor: no hace falta un "
        "sistema perfecto de datos para empezar.' ~1:15.")

# ============================================================
# 11. EMPEZAR CON LO QUE SE TIENE
# ============================================================
s = slide(); base(s)
kicker(s, "Sin excusas")
title(s, [("Se empieza ", SLATE, False), ("con lo que hay", GOLD, False)])
text(s, Inches(0.7), Inches(1.95), Inches(11.6), Inches(0.5),
     [[("Aun con datos y capacidad institucional limitados, toda ciudad puede dar pasos "
        "fundacionales:", 14, RGBColor(0x5C,0x63,0x6E), False, False)]], line_spacing=1.15)
steps = ["Marcos regulatorios y normativas sencillas y claras.",
         "Responsabilidades institucionales y mecanismos de coordinación definidos.",
         "Usar los datos de riesgo disponibles —aun de baja resolución— con registros históricos, "
         "observación y conocimiento local.",
         "Mecanismos prácticos: zonas de no construcción, retiros, control en áreas peligrosas.",
         "Capacitar a funcionarios en uso de datos, herramientas espaciales y monitoreo."]
yy = Inches(2.6)
for i,t in enumerate(steps):
    rect(s, Inches(0.8), yy, Inches(0.5), Inches(0.5), fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, Inches(0.8), yy, Inches(0.5), Inches(0.5),
         [[(str(i+1), 16, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, Inches(1.5), yy, Inches(10.9), Inches(0.6),
         [[(t, 15, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=0)
    yy = yy + Inches(0.78)
footer(s)
notes(s, "Slide profundamente motivacional — el antídoto contra la parálisis. 'No tenemos datos "
        "perfectos' NO es excusa. El conocimiento local vale. La perfección es enemiga de la "
        "acción.\nEnfatizar lo iterativo: se arranca con tamizaje general y, donde el riesgo es "
        "alto, se profundiza (modelación hidráulica, microzonificación sísmica, estabilidad de "
        "laderas).\nTransición: 'A medida que la ciudad madura, la caja de herramientas crece.' ~1:15.")

# ============================================================
# 12. ESCALAR LA AMBICIÓN
# ============================================================
s = slide(); base(s)
kicker(s, "Escalar la ambición")
title(s, [("De herramientas ", SLATE, False), ("básicas", SLATE, True),
          (" a instrumentos ", SLATE, False), ("avanzados", GOLD, True)], size=30)
lw = Inches(5.1)
rect(s, Inches(0.7), Inches(2.3), lw, Inches(2.1), fill=RGBColor(0xEE,0xF1,0xF4),
     line=RGBColor(0xD8,0xDE,0xE5), line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
text(s, Inches(0.95), Inches(2.5), lw-Inches(0.5), Inches(0.5),
     [[("Capacidad inicial", 17, SLATE, True, False)]], space_after=0)
text(s, Inches(0.95), Inches(3.05), lw-Inches(0.5), Inches(1.3),
     [[("Zonificación básica · planificación del uso del suelo · códigos de edificación · "
        "campañas de concientización · incentivos focalizados.", 14, INK, False, False)]], line_spacing=1.25)
rect(s, Inches(5.95), Inches(3.1), Inches(1.0), Inches(0.65), fill=GOLD, shape=MSO_SHAPE.RIGHT_ARROW)
rect(s, Inches(7.1), Inches(2.3), lw, Inches(2.1), fill=RGBColor(0xEF,0xE7,0xD6),
     line=RGBColor(0xE0,0xD4,0xAD), line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
text(s, Inches(7.35), Inches(2.5), lw-Inches(0.5), Inches(0.5),
     [[("Capacidad avanzada", 17, GOLD, True, False)]], space_after=0)
text(s, Inches(7.35), Inches(3.05), lw-Inches(0.5), Inches(1.3),
     [[("Captura de plusvalías (land value capture)", 14, SLATE_DK, True, False),
       (" · finanzas verdes · esquemas de seguros para financiar el desarrollo urbano resiliente.",
        14, INK, False, False)]], line_spacing=1.25)
text(s, Inches(0.7), Inches(4.85), Inches(11.6), Inches(0.6),
     [[("Un proceso multiescala e iterativo: el instrumento se ajusta a la capacidad y al "
        "contexto de cada ciudad.", 14, MUTED, False, False)]], line_spacing=1.2)
footer(s)
notes(s, "Acá conecto con mi propia agenda de investigación: los instrumentos basados en suelo —la "
        "captura de plusvalías— no son un lujo de países ricos; son una vía para FINANCIAR la "
        "resiliencia con el propio valor que crea la ciudad.\nMensaje: no es una escalera que "
        "separa ricos de pobres; es un camino que toda ciudad puede recorrer.\nTransición: 'Todo "
        "esto está ordenado en un manual de tres partes.' ~1:10.")

# ============================================================
# 13. ESTRUCTURA DEL HANDBOOK
# ============================================================
s = slide(); base(s)
kicker(s, "Cómo está organizado")
title(s, [("El Handbook, en tres partes", SLATE, False)])
parts = [("PARTE I · Marco conceptual", RGBColor(0x2F,0x6F,0x7C),
          ["Conceptos, principios, metas y gobernanza","Presenta el marco RCP",
           "La relación entre urbanización y riesgo"]),
         ("PARTE II · Proceso e instrumentos", INFRA,
          ["Cómo incorporar el riesgo al plan","Instrumentos regulatorios y financieros",
           "Monitoreo y evaluación"]),
         ("PARTE III · Medidas RCP por amenaza", ENV,
          ["Amenazas hidrometeorológicas, geológicas y climáticas","Catálogo de ejemplos RCP",
           "Conclusiones"])]
pw = Inches(3.85); px = Inches(0.7); gap = Inches(0.18)
for i,(hd,col,items) in enumerate(parts):
    cx = px + i*(pw+gap)
    rect(s, cx, Inches(2.2), pw, Inches(3.1), fill=WHITE, line=CARD_BD, line_w=Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.04)
    rect(s, cx, Inches(2.2), pw, Inches(0.6), fill=col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    text(s, cx+Inches(0.2), Inches(2.2), pw-Inches(0.4), Inches(0.6),
         [[(hd, 13, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    yy = Inches(3.0)
    for it in items:
        rect(s, cx+Inches(0.25), yy+Inches(0.09), Inches(0.1), Inches(0.1), fill=col, shape=MSO_SHAPE.OVAL)
        text(s, cx+Inches(0.5), yy, pw-Inches(0.7), Inches(0.7),
             [[(it, 12.5, INK, False, False)]], line_spacing=1.15, space_after=0)
        yy = yy + Inches(0.72)
source(s, "Un marco integral para transformar la urbanización de motor de riesgo en catalizador de "
          "resiliencia y capacidad.")
footer(s)
notes(s, "Recorrido rápido — dar el mapa, no leer todo. 'Parte I: el por qué y el marco. Parte II: "
        "el cómo operativo. Parte III: el recetario por tipo de amenaza.'\nRecalcar que es una GUÍA "
        "PRÁCTICA, no un compendio exhaustivo: pensada para usarse. ~0:50.")

# ============================================================
# 14. ¿PARA QUIÉN?
# ============================================================
s = slide(); base(s)
kicker(s, "¿A quién interpela?")
title(s, [("Una herramienta para quienes ", SLATE, False),
          ("deciden y construyen", GOLD, False), (" ciudad", SLATE, False)], size=30)
aud = [("Gobiernos y tomadores de decisión","Nacionales y locales"),
       ("Planificadores y técnicos","Urbanismo, DRR y resiliencia climática"),
       ("Investigadores y academia","Evidencia para la política"),
       ("ONG y organizaciones de base","Adaptación basada en ecosistemas"),
       ("Sector privado y financiero","Inversión resiliente"),
       ("Comunidades","Justicia climática e inclusión")]
cw = Inches(3.85); ch = Inches(1.35); cx0 = Inches(0.7); cy0 = Inches(2.2)
gx = Inches(0.18); gy = Inches(0.25)
for i,(h,sub) in enumerate(aud):
    cx = cx0 + (i%3)*(cw+gx); cy = cy0 + (i//3)*(ch+gy)
    rect(s, cx, cy, cw, ch, fill=WHITE, line=CARD_BD, line_w=Pt(1), shadow=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05)
    rect(s, cx, cy, cw, Pt(3), fill=GOLD)
    text(s, cx+Inches(0.25), cy+Inches(0.22), cw-Inches(0.5), Inches(0.7),
         [[(h, 15, SLATE_DK, True, False)]], line_spacing=1.05, space_after=0)
    text(s, cx+Inches(0.25), cy+Inches(0.9), cw-Inches(0.5), Inches(0.4),
         [[(sub, 12, MUTED, False, False)]], space_after=0)
footer(s)
notes(s, "La resiliencia es deporte de equipo. Nadie la logra solo. Mirar al público y nombrar "
        "quiénes están en la sala (autoridades, colegas, estudiantes) — 'cada uno de ustedes tiene "
        "una palanca'.\nSubir la energía hacia el cierre. ~0:50.")

# ============================================================
# 15. CIERRE / LLAMADO A LA ACCIÓN (oscura)
# ============================================================
s = slide(); base(s, dark=True)
kicker(s, "El llamado", color=GOLD_LT)
rect(s, Inches(0.7), Inches(1.9), Pt(5), Inches(2.2), fill=GOLD)
text(s, Inches(1.0), Inches(1.9), Inches(10.5), Inches(2.3),
     [[("La urbanización puede pasar de ser un ", 30, WHITE, True, False),
       ("motor de riesgo", 30, RGBColor(0xE7,0xB8,0xB8), True, False),
       (" a un ", 30, WHITE, True, False),
       ("catalizador de resiliencia", 30, GOLD_LT, True, False),
       (".", 30, WHITE, True, False)]], line_spacing=1.2)
text(s, Inches(1.0), Inches(4.6), Inches(10.5), Inches(1.0),
     [[("No es un problema de destino, sino de ", 18, RGBColor(0xCF,0xD5,0xDD), False, False),
       ("decisiones de planificación", 18, WHITE, True, False),
       (". Y esas decisiones se toman ", 18, RGBColor(0xCF,0xD5,0xDD), False, False),
       ("hoy", 18, GOLD_LT, True, False),
       (".", 18, RGBColor(0xCF,0xD5,0xDD), False, False)]], line_spacing=1.25)
text(s, Inches(1.0), Inches(5.7), Inches(10), Inches(0.7),
     [[("Planifiquemos ciudades que protejan a su gente.", 22, WHITE, True, False)]], space_after=0)
notes(s, "Cierre con convicción y calma. Bajar el ritmo, contacto visual. Repetir la idea-ancla: "
        "la ciudad como solución.\nÚltima frase con énfasis: 'El riesgo se acumula en silencio; la "
        "resiliencia se construye con decisiones. Empecemos hoy.' Pausa. Gracias. ~1:00.")

# ============================================================
# 16. GRACIAS
# ============================================================
s = slide(); base(s)
kicker(s, "Gracias")
title(s, [("Conversemos", SLATE, False)])
text(s, Inches(0.7), Inches(2.3), Inches(11), Inches(0.6),
     [[("Prof. Cynthia Goytia (PhD)", 20, SLATE_DK, True, False)],
      [("Economía Urbana · Universidad Torcuato Di Tella · CIPUV", 14, MUTED, False, False)]],
     line_spacing=1.2)
text(s, Inches(0.7), Inches(3.5), Inches(11), Inches(0.5),
     [[("✉  cgoytia@utdt.edu      🌐  cgoytia-droid.github.io", 16, SLATE, False, False)]], space_after=0)
text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.3),
     [[("Basado en: World Bank & GFDRR (2025). Handbook for Livable and Resilient Cities: "
        "Integrating Hazard and Risk Information into Urban Planning. Washington, DC: World Bank.",
        11, MUTED, False, True)],
      [("Datos: World Bank 2022a; United Nations 2018; CRED & UNDRR 2019; Winsemius et al. 2016; "
        "Our World in Data 2025.", 11, MUTED, False, True)]], line_spacing=1.2)
footer(s)
notes(s, "Abrir el diálogo con una pregunta al público: '¿Qué instrumento RCP creen que su ciudad "
        "podría aplicar mañana?'\nTener a mano ejemplos locales de captura de plusvalías y de zonas "
        "de no construcción para responder. Agradecer a los organizadores.")

out = "/home/user/cgoytia-droid.github.io/presentations/ciudades-habitables-resilientes.pptx"
prs.save(out)
print("OK:", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
